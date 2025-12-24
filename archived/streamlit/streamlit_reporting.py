"""
PCA Agent - Automated Reporting Module

Upload report templates (XLSX, CSV, PPTX) and automatically populate them with campaign data.
The system understands template structure and maps data to the correct locations.

Run with:
    streamlit run streamlit_reporting.py --server.port 8504
"""

import os
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
from datetime import datetime
import io
import pickle
import hashlib
import logging

import streamlit as st
import pandas as pd
import numpy as np
from dotenv import load_dotenv

# Add project root to path
current_dir = os.path.dirname(os.path.abspath(__file__))
if current_dir not in sys.path:
    sys.path.insert(0, current_dir)

load_dotenv()

# =============================================================================
# LOGGING SETUP
# =============================================================================

# Create logs directory
LOG_DIR = Path(current_dir) / "logs"
LOG_DIR.mkdir(exist_ok=True)

# Configure logging
log_file = LOG_DIR / f"reporting_{datetime.now().strftime('%Y-%m-%d')}.log"

# Create formatter
log_formatter = logging.Formatter(
    '%(asctime)s | %(levelname)-8s | %(funcName)s:%(lineno)d | %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)

# File handler
file_handler = logging.FileHandler(log_file, encoding='utf-8')
file_handler.setFormatter(log_formatter)
file_handler.setLevel(logging.DEBUG)

# Console handler (for Streamlit terminal)
console_handler = logging.StreamHandler()
console_handler.setFormatter(log_formatter)
console_handler.setLevel(logging.INFO)

# Create logger
logger = logging.getLogger('reporting')
logger.setLevel(logging.DEBUG)
logger.handlers = []  # Clear existing handlers
logger.addHandler(file_handler)
logger.addHandler(console_handler)

logger.info("=" * 60)
logger.info("REPORTING MODULE STARTED")
logger.info("=" * 60)

# Cache directory for persisting uploaded files
CACHE_DIR = Path(current_dir) / ".reporting_cache"
CACHE_DIR.mkdir(exist_ok=True)

# Page config
st.set_page_config(
    page_title="PCA Agent - Reporting",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

# Inject CSS
st.markdown("""
<style>
    .main-header {
        font-size: 2.4rem;
        font-weight: 700;
        color: #5b6ef5;
        margin-bottom: 1rem;
    }
    .section-header {
        font-size: 1.6rem;
        font-weight: 600;
        color: #667eea;
        margin-top: 1.5rem;
        margin-bottom: 0.5rem;
    }
    .info-box {
        background-color: rgba(91, 110, 245, 0.1);
        border-left: 4px solid #5b6ef5;
        padding: 1rem;
        margin: 1rem 0;
        border-radius: 4px;
    }
</style>
""", unsafe_allow_html=True)


def save_to_cache(key: str, data: Any, filename: str = None):
    """Save data to persistent cache."""
    try:
        cache_data = {
            'data': data,
            'filename': filename,
            'timestamp': datetime.now().isoformat()
        }
        cache_path = CACHE_DIR / f"{key}.pkl"
        with open(cache_path, 'wb') as f:
            pickle.dump(cache_data, f)
        logger.debug(f"Cache saved: {key} (filename: {filename})")
    except Exception as e:
        logger.error(f"Cache save error for {key}: {e}")


def load_from_cache(key: str) -> Tuple[Any, Optional[str]]:
    """Load data from persistent cache. Returns (data, filename) or (None, None)."""
    try:
        cache_path = CACHE_DIR / f"{key}.pkl"
        if cache_path.exists():
            with open(cache_path, 'rb') as f:
                cache_data = pickle.load(f)
                logger.debug(f"Cache loaded: {key} (filename: {cache_data.get('filename')})")
                return cache_data.get('data'), cache_data.get('filename')
    except Exception as e:
        logger.error(f"Cache load error for {key}: {e}")
    return None, None


def clear_cache():
    """Clear all cached files."""
    try:
        for cache_file in CACHE_DIR.glob("*.pkl"):
            cache_file.unlink()
        logger.info("Cache cleared")
    except Exception as e:
        logger.error(f"Cache clear error: {e}")


def init_session_state():
    """Initialize session state variables, loading from cache if available."""
    defaults = {
        "template_file": None,
        "template_type": None,
        "template_structure": None,
        "data_file": None,
        "data_df": None,
        "mapping_config": {},
        "generated_report": None,
        "template_filename": None,
        "data_filename": None,
        "cache_loaded": False,
    }
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value
    
    # Load from cache on first run (only once per session)
    if not st.session_state.cache_loaded:
        st.session_state.cache_loaded = True
        
        # Load cached template (stored as bytes, convert to BytesIO)
        cached_template_bytes, template_filename = load_from_cache("template_file")
        if cached_template_bytes is not None:
            st.session_state.template_file = io.BytesIO(cached_template_bytes)
            st.session_state.template_filename = template_filename
            if template_filename:
                st.session_state.template_type = template_filename.split('.')[-1].lower()
        
        # Load cached template structure
        cached_structure, _ = load_from_cache("template_structure")
        if cached_structure is not None:
            st.session_state.template_structure = cached_structure
        
        # Load cached data
        cached_df, data_filename = load_from_cache("data_df")
        if cached_df is not None:
            st.session_state.data_df = cached_df
            st.session_state.data_filename = data_filename
            st.session_state.data_file = data_filename
        
        # Load cached mapping config
        cached_mapping, _ = load_from_cache("mapping_config")
        if cached_mapping is not None:
            st.session_state.mapping_config = cached_mapping


def render_sidebar() -> str:
    """Render sidebar navigation."""
    with st.sidebar:
        st.markdown("## 📊 Reporting Module")
        st.markdown("*Automated Report Generation*")
        
        st.info("""
        Upload your source data and report template to automatically generate 
        populated reports in the same format.
        """)
        
        st.divider()
        
        page = st.radio(
            "Navigation",
            options=["Upload Data & Template", "Data Mapping", "Generate Report", "Settings"],
            index=0,
        )
        
        st.divider()
        
        # Status indicators
        st.subheader("Status")
        
        if st.session_state.template_file:
            template_name = st.session_state.get('template_filename') or st.session_state.template_type
            st.success(f"✅ Template: {template_name}")
        else:
            st.warning("⚠️ No template uploaded")
        
        if st.session_state.data_df is not None:
            data_name = st.session_state.get('data_filename') or f"{len(st.session_state.data_df)} rows"
            st.success(f"✅ Data: {data_name}")
        else:
            st.warning("⚠️ No data loaded")
        
        if st.session_state.generated_report:
            st.success("✅ Report generated")
        
        st.divider()
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("🔄 Reset All"):
                for key in ["template_file", "template_type", "template_structure", 
                           "data_file", "data_df", "mapping_config", "generated_report",
                           "template_filename", "data_filename"]:
                    st.session_state[key] = None if key != "mapping_config" else {}
                clear_cache()
                st.rerun()
        
        with col2:
            if st.button("🗑️ Clear Cache"):
                clear_cache()
                st.toast("Cache cleared!", icon="✅")
    
    return page.lower().replace(" ", "_")


def render_upload_page():
    """Render combined data and template upload page."""
    st.markdown('<div class="main-header">📤 Upload Data & Template</div>', unsafe_allow_html=True)
    
    st.markdown("""
    ### Step 1: Upload Source Data
    Upload your campaign performance data that will populate the report template.
    """)
    
    # Data upload section
    col1, col2 = st.columns([2, 1])
    
    with col1:
        data_file = st.file_uploader(
            "📊 Upload Source Data",
            type=["csv", "xlsx", "xls"],
            help="Upload your campaign performance data (CSV or Excel)",
            key="data_uploader"
        )
    
    with col2:
        st.markdown("#### Supported Formats")
        st.markdown("- CSV (.csv)")
        st.markdown("- Excel (.xlsx, .xls)")
    
    # Show cached data info if available and no new file uploaded
    if data_file is None and st.session_state.data_df is not None:
        cached_name = st.session_state.get('data_filename', 'cached data')
        st.info(f"📁 **Using cached data**: {cached_name} ({len(st.session_state.data_df)} rows)")
    
    if data_file is not None:
        try:
            # Load data
            if data_file.name.endswith('.csv'):
                df = pd.read_csv(data_file)
            else:
                df = pd.read_excel(data_file)
            
            # Try to normalize data (handle messy formats)
            try:
                from src.utils.data_normalizer import normalize_dataframe, validate_data_quality, detect_column_type
                df, column_types = normalize_dataframe(df, auto_detect=True)
                quality_report = validate_data_quality(df)
                st.session_state.column_types = column_types
                st.session_state.quality_report = quality_report
            except ImportError:
                column_types = {}
                quality_report = None
            
            st.session_state.data_df = df
            st.session_state.data_file = data_file.name
            st.session_state.data_filename = data_file.name
            
            # Save to cache
            save_to_cache("data_df", df, data_file.name)
            
            st.success(f"✅ Data loaded: {len(df)} rows, {len(df.columns)} columns")
            
            # Show data quality warnings if any
            if quality_report and quality_report.get('warnings'):
                with st.expander("⚠️ Data Quality Warnings", expanded=True):
                    for warning in quality_report['warnings']:
                        st.warning(warning)
            
            # Data preview
            with st.expander("📋 Data Preview", expanded=True):
                col1, col2, col3, col4 = st.columns(4)
                with col1:
                    st.metric("Rows", len(df))
                with col2:
                    st.metric("Columns", len(df.columns))
                with col3:
                    numeric_cols = len(df.select_dtypes(include=['number']).columns)
                    st.metric("Numeric Cols", numeric_cols)
                with col4:
                    text_cols = len(df.select_dtypes(include=['object']).columns)
                    st.metric("Text Cols", text_cols)
                
                st.dataframe(df.head(10), use_container_width=True)
                
                # Column info with detected types
                with st.expander("ℹ️ Column Information"):
                    col_info_data = {
                        'Column': df.columns,
                        'Data Type': df.dtypes.astype(str),
                        'Detected Type': [column_types.get(col, 'unknown') for col in df.columns] if column_types else ['N/A'] * len(df.columns),
                        'Non-Null': df.count().values,
                        'Null %': [f"{df[col].isna().mean():.1%}" for col in df.columns],
                        'Unique': [df[col].nunique() for col in df.columns]
                    }
                    col_info = pd.DataFrame(col_info_data)
                    st.dataframe(col_info, use_container_width=True)
        
        except Exception as e:
            st.error(f"❌ Error loading data: {str(e)}")
    
    st.divider()
    
    # Template upload section
    st.markdown("""
    ### Step 2: Upload Report Template
    
    Upload your report template with placeholders. Supported formats:
    - **Excel (.xlsx, .xls)** - Spreadsheet templates with formulas and formatting
    - **CSV (.csv)** - Simple tabular templates
    - **PowerPoint (.pptx)** - Presentation templates with charts and tables
    
    The system will analyze the template structure and identify data placeholders.
    """)
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        template_file = st.file_uploader(
            "📄 Upload Report Template",
            type=["xlsx", "xls", "csv", "pptx"],
            help="Upload your report template in Excel, CSV, or PowerPoint format",
            key="template_uploader"
        )
    
    with col2:
        st.markdown("#### Placeholder Formats")
        st.markdown("- `{{field}}`")
        st.markdown("- `{field}`")
        st.markdown("- `[field]`")
        st.markdown("- `<field>`")
    
    # Show cached template info if available and no new file uploaded
    if template_file is None and st.session_state.template_file is not None:
        cached_name = st.session_state.get('template_filename', 'cached template')
        st.info(f"📁 **Using cached template**: {cached_name}")
    
    if template_file is not None:
        file_ext = template_file.name.split('.')[-1].lower()
        
        # Read file bytes for caching (since UploadedFile can't be pickled directly)
        template_bytes = template_file.read()
        template_file.seek(0)  # Reset for later use
        
        st.session_state.template_file = io.BytesIO(template_bytes)
        st.session_state.template_type = file_ext
        st.session_state.template_filename = template_file.name
        
        # Save to cache
        save_to_cache("template_file", template_bytes, template_file.name)
        
        st.success(f"✅ Template uploaded: {template_file.name}")
        
        # Analyze template
        with st.spinner("🔍 Analyzing template structure..."):
            template_for_analysis = io.BytesIO(template_bytes)
            structure = analyze_template(template_for_analysis, file_ext)
            st.session_state.template_structure = structure
            
            # Cache the structure too
            save_to_cache("template_structure", structure)
        
        # Display template analysis
        st.markdown('<div class="section-header">Template Analysis</div>', unsafe_allow_html=True)
        
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.metric("Format", file_ext.upper())
        
        with col2:
            if file_ext in ['xlsx', 'xls', 'csv']:
                st.metric("Sheets/Tables", structure.get('sheet_count', 1))
        
        with col3:
            st.metric("Data Fields", structure.get('field_count', 0))
        
        with col4:
            mode = structure.get('mode', 'placeholder')
            mode_label = "📝 Placeholder" if mode == 'placeholder' else "📊 Column"
            st.metric("Mode", mode_label)
        
        # Show mode explanation
        mode = structure.get('mode', 'placeholder')
        if mode == 'column':
            st.info("📊 **Column Mode Detected**: Your template has column headers. Data rows from your input file will be populated into matching columns.")
        elif structure.get('placeholders'):
            st.info("📝 **Placeholder Mode Detected**: Your template has placeholders like `{{field}}`. These will be replaced with aggregated values.")
        
        # Show detected placeholders
        if structure.get('placeholders'):
            with st.expander("🔍 Detected Placeholders", expanded=True):
                st.markdown("These fields were detected in your template:")
                for placeholder in structure['placeholders']:
                    st.markdown(f"- `{placeholder['name']}` ({placeholder['type']}) - {placeholder['location']}")
        
        # Show detected columns (for column mode)
        if mode == 'column' and structure.get('template_columns'):
            with st.expander("📋 Detected Template Columns", expanded=True):
                st.markdown("These column headers were found in your template:")
                for col in structure['template_columns']:
                    st.markdown(f"- `{col['name']}` (Sheet: {col['sheet']}, Column: {col.get('column_letter', 'N/A')})")
        
        # Show preview
        if file_ext in ['xlsx', 'xls', 'csv']:
            with st.expander("📋 Template Preview", expanded=False):
                preview_df = load_template_preview(template_file, file_ext)
                if preview_df is not None:
                    st.dataframe(preview_df.head(10), use_container_width=True)


def analyze_template(file, file_type: str) -> Dict[str, Any]:
    """
    Analyze template structure and identify placeholders OR column headers.
    
    Supports two modes:
    1. Placeholder mode: Template has {{field}}, {field}, [field], <field> placeholders
    2. Column mode: Template has column headers that match data columns (for data population)
    
    Args:
        file: Uploaded file object
        file_type: File extension (xlsx, csv, pptx)
    
    Returns:
        Dictionary with template structure information
    """
    logger.info(f"Analyzing template (type: {file_type})")
    
    structure = {
        'file_type': file_type,
        'placeholders': [],
        'field_count': 0,
        'sheet_count': 0,
        'template_columns': [],  # Column headers found in template
        'data_start_row': 2,     # Row where data should start (after headers)
        'mode': 'placeholder'    # 'placeholder' or 'column'
    }
    
    try:
        if file_type in ['xlsx', 'xls']:
            import openpyxl
            
            # Load workbook
            wb = openpyxl.load_workbook(file, data_only=False)
            structure['sheet_count'] = len(wb.sheetnames)
            logger.debug(f"Loaded workbook with {len(wb.sheetnames)} sheets: {wb.sheetnames}")
            
            placeholders = []
            template_columns = []
            
            # Scan each sheet for placeholders AND column headers
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                # Find header row (row with most non-empty cells in first 20 rows)
                header_row = 1
                max_filled = 0
                for row_idx in range(1, min(21, sheet.max_row + 1)):
                    filled = sum(1 for cell in sheet[row_idx] if cell.value is not None)
                    if filled > max_filled:
                        max_filled = filled
                        header_row = row_idx
                
                # Extract column headers from detected header row
                from openpyxl.utils import get_column_letter as gcl
                header_row_values = []
                for cell in sheet[header_row]:
                    # Skip merged cells
                    if not hasattr(cell, 'column') or cell.column is None:
                        continue
                    if cell.value:
                        header_row_values.append({
                            'name': str(cell.value),
                            'column_letter': gcl(cell.column),
                            'sheet': sheet_name,
                            'header_row': header_row
                        })
                
                if header_row_values:
                    template_columns.extend(header_row_values)
                    # Store the data start row for this sheet
                    structure[f'data_start_row_{sheet_name}'] = header_row + 1
                
                # Also scan for explicit placeholders (but NOT formulas)
                import re
                placeholder_pattern = re.compile(r'\{\{[^}]+\}\}|\[\[[^\]]+\]\]|<<[^>]+>>')
                
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            cell_str = str(cell.value)
                            # Skip formulas
                            if cell_str.startswith('='):
                                continue
                            # Only detect explicit placeholder patterns: {{field}}, [[field]], <<field>>
                            if placeholder_pattern.search(cell_str):
                                placeholders.append({
                                    'name': cell.value,
                                    'type': 'text',
                                    'location': f"{sheet_name}!{cell.coordinate}",
                                    'sheet': sheet_name,
                                    'cell': cell.coordinate
                                })
            
            structure['placeholders'] = placeholders
            structure['template_columns'] = template_columns
            structure['field_count'] = len(placeholders) if placeholders else len(template_columns)
            
            # Determine mode based on what was found
            if placeholders:
                structure['mode'] = 'placeholder'
                logger.info(f"Template mode: PLACEHOLDER ({len(placeholders)} placeholders found)")
                for p in placeholders[:5]:
                    logger.debug(f"  Placeholder: {p['name']} at {p['location']}")
            elif template_columns:
                structure['mode'] = 'column'
                logger.info(f"Template mode: COLUMN ({len(template_columns)} columns found)")
                for c in template_columns[:10]:
                    logger.debug(f"  Column: {c['name']} ({c['column_letter']}) in {c['sheet']}")
            
        elif file_type == 'csv':
            df = pd.read_csv(file)
            structure['sheet_count'] = 1
            
            # Store all columns as template columns
            template_columns = [{'name': col, 'column_letter': None, 'sheet': 'main'} for col in df.columns]
            structure['template_columns'] = template_columns
            
            # Check for placeholder columns
            placeholders = []
            for col in df.columns:
                if any(pattern in str(col) for pattern in ['{{', '{', '[', '<']):
                    placeholders.append({
                        'name': col,
                        'type': 'column',
                        'location': f"Column: {col}",
                        'sheet': 'main',
                        'cell': col
                    })
            
            structure['placeholders'] = placeholders
            structure['field_count'] = len(placeholders) if placeholders else len(template_columns)
            structure['mode'] = 'placeholder' if placeholders else 'column'
            
        elif file_type == 'pptx':
            from pptx import Presentation
            
            prs = Presentation(file)
            structure['sheet_count'] = len(prs.slides)
            
            placeholders = []
            
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        if any(pattern in shape.text for pattern in ['{{', '{', '[', '<']):
                            placeholders.append({
                                'name': shape.text,
                                'type': 'text',
                                'location': f"Slide {slide_idx + 1}",
                                'sheet': f"slide_{slide_idx}",
                                'cell': shape.name
                            })
            
            structure['placeholders'] = placeholders
            structure['field_count'] = len(placeholders)
    
    except Exception as e:
        st.error(f"Error analyzing template: {str(e)}")
    
    return structure


def load_template_preview(file, file_type: str) -> Optional[pd.DataFrame]:
    """Load template preview as DataFrame."""
    try:
        if file_type in ['xlsx', 'xls']:
            return pd.read_excel(file, nrows=10)
        elif file_type == 'csv':
            return pd.read_csv(file, nrows=10)
    except Exception as e:
        st.warning(f"Could not load preview: {str(e)}")
    return None


def render_data_mapping_page():
    """Render data mapping configuration page."""
    st.markdown('<div class="main-header">🔗 Data Mapping</div>', unsafe_allow_html=True)
    
    if not st.session_state.template_file:
        st.warning("⚠️ Please upload a template first in the 'Upload Data & Template' page")
        return
    
    if st.session_state.data_df is None:
        st.warning("⚠️ Please upload source data first in the 'Upload Data & Template' page")
        return
    
    st.markdown("""
    ### Configure Field Mapping
    
    Map template placeholders to your data columns. The system has auto-suggested matches based on field names.
    """)
    
    st.divider()
    
    df = st.session_state.data_df
    
    # Show data summary
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("Data Source", st.session_state.data_file)
    with col2:
        st.metric("Rows", len(df))
    with col3:
        st.metric("Columns", len(df.columns))
    
    st.divider()
    
    # Mapping configuration
    st.markdown('<div class="section-header">Field Mappings</div>', unsafe_allow_html=True)
    
    template_structure = st.session_state.template_structure
    mode = template_structure.get('mode', 'placeholder') if template_structure else 'placeholder'
    
    if mode == 'placeholder' and template_structure.get('placeholders'):
        # Original placeholder mode
        st.markdown("Map template placeholders to data columns:")
        
        mapping_config = {}
        
        for placeholder in template_structure['placeholders']:
            col1, col2, col3 = st.columns([2, 2, 1])
            
            with col1:
                st.text_input(
                    "Template Field",
                    value=placeholder['name'],
                    disabled=True,
                    key=f"template_{placeholder['location']}"
                )
            
            with col2:
                # Auto-suggest matching column
                suggested_col = find_matching_column(placeholder['name'], df.columns)
                
                data_column = st.selectbox(
                    "Data Column",
                    options=['None'] + list(df.columns),
                    index=list(df.columns).index(suggested_col) + 1 if suggested_col else 0,
                    key=f"mapping_{placeholder['location']}"
                )
                
                if data_column != 'None':
                    mapping_config[placeholder['name']] = {
                        'data_column': data_column,
                        'location': placeholder['location'],
                        'type': placeholder['type']
                    }
            
            with col3:
                st.caption(placeholder['location'])
        
        st.session_state.mapping_config = mapping_config
        
        if st.button("💾 Save Mapping", type="primary"):
            save_to_cache("mapping_config", mapping_config)
            st.success("✅ Mapping configuration saved!")
    
    elif mode == 'column' and template_structure.get('template_columns'):
        # Column-based data population mode
        st.info("📋 **Column Mode**: Template columns will be matched to data columns. Data rows will be populated into the template.")
        
        template_cols = template_structure['template_columns']
        data_cols = list(df.columns)
        
        st.markdown("### Auto-Matched Columns")
        
        mapping_config = {'_mode': 'column', '_column_mappings': {}}
        matched_count = 0
        
        for tcol in template_cols:
            col1, col2, col3 = st.columns([2, 2, 1])
            
            with col1:
                st.text_input(
                    "Template Column",
                    value=tcol['name'],
                    disabled=True,
                    key=f"tcol_{tcol['sheet']}_{tcol['name']}"
                )
            
            with col2:
                # Auto-match by name (case-insensitive, flexible matching)
                suggested = find_matching_column(tcol['name'], data_cols)
                
                data_column = st.selectbox(
                    "Data Column",
                    options=['(skip)'] + data_cols,
                    index=data_cols.index(suggested) + 1 if suggested else 0,
                    key=f"dcol_{tcol['sheet']}_{tcol['name']}"
                )
                
                if data_column != '(skip)':
                    mapping_config['_column_mappings'][tcol['name']] = {
                        'data_column': data_column,
                        'column_letter': tcol.get('column_letter'),
                        'sheet': tcol['sheet'],
                        'header_row': tcol.get('header_row', 1)
                    }
                    matched_count += 1
            
            with col3:
                if suggested:
                    st.success("✓ Auto")
                else:
                    st.warning("Manual")
        
        st.divider()
        st.metric("Matched Columns", f"{matched_count} / {len(template_cols)}")
        
        st.session_state.mapping_config = mapping_config
        
        if st.button("💾 Save Column Mapping", type="primary"):
            save_to_cache("mapping_config", mapping_config)
            st.success(f"✅ Column mapping saved! {matched_count} columns will be populated.")
    
    else:
        st.info("No placeholders or column headers detected in template. The entire data will be inserted as a new sheet.")


def find_matching_column(placeholder: str, columns: List[str]) -> Optional[str]:
    """
    Find best matching column for a placeholder using fuzzy matching.
    
    Handles:
    - Different naming conventions (snake_case, camelCase, spaces)
    - Common abbreviations (impr -> impressions, conv -> conversions)
    - Typos (via fuzzy matching)
    - Column aliases (spend vs cost, revenue vs sales)
    """
    try:
        from src.utils.data_normalizer import find_best_match, normalize_column_name
        
        # Use advanced fuzzy matching
        match, score = find_best_match(placeholder, list(columns), threshold=0.5)
        if match:
            return match
    except ImportError:
        pass
    
    # Fallback to basic matching
    clean_placeholder = placeholder.strip('{}[]<>').lower().replace('_', ' ').replace('-', ' ')
    
    # Try exact match
    for col in columns:
        if col.lower() == clean_placeholder:
            return col
    
    # Try partial match
    for col in columns:
        col_clean = col.lower().replace('_', ' ').replace('-', ' ')
        if clean_placeholder in col_clean or col_clean in clean_placeholder:
            return col
    
    # Try common aliases
    aliases = {
        'spend': ['cost', 'amount', 'budget'],
        'revenue': ['sales', 'income', 'value'],
        'impressions': ['impr', 'views', 'imp'],
        'clicks': ['click', 'clk'],
        'conversions': ['conv', 'converts', 'results'],
        'ctr': ['click through rate', 'clickthrough'],
        'cpc': ['cost per click'],
        'cpa': ['cost per acquisition', 'cost per conversion'],
        'roas': ['return on ad spend', 'roi'],
    }
    
    for canonical, alias_list in aliases.items():
        if canonical in clean_placeholder or clean_placeholder in canonical:
            for col in columns:
                col_clean = col.lower().replace('_', ' ').replace('-', ' ')
                if canonical in col_clean:
                    return col
        for alias in alias_list:
            if alias in clean_placeholder or clean_placeholder in alias:
                for col in columns:
                    col_clean = col.lower().replace('_', ' ').replace('-', ' ')
                    if canonical in col_clean or alias in col_clean:
                        return col
    
    return None


def apply_date_grouping(df: pd.DataFrame, grouping: str, week_start: str = "Monday") -> pd.DataFrame:
    """
    Apply date grouping to aggregate data by week or month.
    
    Args:
        df: DataFrame with a date column
        grouping: "Weekly (Week Start)" or "Monthly (Month Start)"
        week_start: "Monday" or "Sunday"
    
    Returns:
        Aggregated DataFrame with dates converted to period start dates
    """
    # Find date column
    date_col = None
    for col in df.columns:
        col_lower = col.lower()
        if col_lower in ['date', 'day', 'report_date', 'campaign_date']:
            date_col = col
            break
        if 'date' in col_lower:
            date_col = col
            break
    
    if date_col is None:
        return df  # No date column found, return as-is
    
    # Make a copy
    result_df = df.copy()
    
    # Convert to datetime
    result_df[date_col] = pd.to_datetime(result_df[date_col], errors='coerce')
    
    # Identify numeric columns for aggregation
    numeric_cols = result_df.select_dtypes(include=['number']).columns.tolist()
    
    # Identify grouping columns (non-numeric, non-date)
    group_cols = []
    for col in result_df.columns:
        if col == date_col:
            continue
        if col not in numeric_cols:
            # Check if it's a categorical/text column worth grouping by
            if result_df[col].nunique() < 50:  # Reasonable number of unique values
                group_cols.append(col)
    
    if "Weekly" in grouping:
        # Convert to week start date
        # Monday = 0, Sunday = 6
        week_start_offset = 0 if week_start == "Monday" else 6
        
        # Get the start of the week for each date
        result_df['_period_start'] = result_df[date_col].apply(
            lambda x: x - pd.Timedelta(days=(x.weekday() - week_start_offset) % 7) if pd.notna(x) else x
        )
        
    elif "Monthly" in grouping:
        # Convert to month start date
        result_df['_period_start'] = result_df[date_col].apply(
            lambda x: x.replace(day=1) if pd.notna(x) else x
        )
    else:
        return df  # No grouping needed
    
    # Replace date column with period start
    result_df[date_col] = result_df['_period_start']
    result_df = result_df.drop(columns=['_period_start'])
    
    # Aggregate by period and group columns
    agg_cols = [date_col] + group_cols
    
    # Build aggregation dict
    agg_dict = {}
    for col in numeric_cols:
        agg_dict[col] = 'sum'
    
    # Group and aggregate
    if agg_dict:
        result_df = result_df.groupby(agg_cols, as_index=False).agg(agg_dict)
    
    # Sort by date
    result_df = result_df.sort_values(date_col)
    
    return result_df


def render_generate_report_page():
    """Render report generation page."""
    st.markdown('<div class="main-header">📊 Generate Report</div>', unsafe_allow_html=True)
    
    if not st.session_state.template_file:
        st.warning("⚠️ Please upload a template first")
        return
    
    if st.session_state.data_df is None:
        st.warning("⚠️ Please upload data and configure mapping")
        return
    
    st.markdown("""
    ### Generate Your Report
    
    Review the configuration and generate your populated report.
    """)
    
    st.divider()
    
    # Configuration summary
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.metric("Template", st.session_state.template_type.upper())
    
    with col2:
        st.metric("Data Rows", len(st.session_state.data_df))
    
    with col3:
        st.metric("Mapped Fields", len(st.session_state.mapping_config))
    
    st.divider()
    
    # Generation options
    st.markdown('<div class="section-header">Generation Options</div>', unsafe_allow_html=True)
    
    col1, col2 = st.columns(2)
    
    with col1:
        report_name = st.text_input(
            "Report Name",
            value=f"Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            help="Name for the generated report"
        )
    
    with col2:
        aggregation_method = st.selectbox(
            "Data Aggregation",
            options=["Sum", "Average", "Latest", "All Rows"],
            help="How to aggregate data if multiple rows match"
        )
    
    col3, col4 = st.columns(2)
    
    with col3:
        date_grouping = st.selectbox(
            "Date Grouping",
            options=["Daily (No Change)", "Weekly (Week Start)", "Monthly (Month Start)"],
            help="Aggregate data by week or month. Dates will be converted to period start dates."
        )
    
    with col4:
        week_start_day = st.selectbox(
            "Week Starts On",
            options=["Monday", "Sunday"],
            help="Which day should be considered the start of the week",
            disabled=date_grouping != "Weekly (Week Start)"
        )
    
    include_charts = st.checkbox("Include Charts", value=True, help="Generate charts from data")
    include_summary = st.checkbox("Include Summary Sheet", value=True, help="Add executive summary")
    
    st.divider()
    
    # Generate button
    if st.button("🚀 Generate Report", type="primary", use_container_width=True):
        with st.spinner("📊 Generating report..."):
            try:
                # Apply date grouping if needed
                working_df = st.session_state.data_df.copy()
                
                if date_grouping != "Daily (No Change)":
                    working_df = apply_date_grouping(
                        working_df, 
                        date_grouping, 
                        week_start_day
                    )
                    st.info(f"📅 Data grouped by {date_grouping.split('(')[0].strip()}: {len(working_df)} rows")
                
                generated_file = generate_report(
                    template_file=st.session_state.template_file,
                    template_type=st.session_state.template_type,
                    data_df=working_df,
                    mapping_config=st.session_state.mapping_config,
                    report_name=report_name,
                    aggregation=aggregation_method.lower(),
                    include_charts=include_charts,
                    include_summary=include_summary
                )
                
                st.session_state.generated_report = generated_file
                st.success("✅ Report generated successfully!")
                
            except Exception as e:
                logger.error(f"Error generating report: {str(e)}", exc_info=True)
                st.error(f"❌ Error generating report: {str(e)}")
    
    # Download generated report
    if st.session_state.generated_report:
        st.divider()
        st.markdown('<div class="section-header">Download Report</div>', unsafe_allow_html=True)
        
        file_ext = st.session_state.template_type
        mime_types = {
            'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'xls': 'application/vnd.ms-excel',
            'csv': 'text/csv',
            'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        }
        
        st.download_button(
            label=f"📥 Download {report_name}.{file_ext}",
            data=st.session_state.generated_report,
            file_name=f"{report_name}.{file_ext}",
            mime=mime_types.get(file_ext, 'application/octet-stream'),
            use_container_width=True
        )


def generate_report(
    template_file,
    template_type: str,
    data_df: pd.DataFrame,
    mapping_config: Dict,
    report_name: str,
    aggregation: str = 'sum',
    include_charts: bool = True,
    include_summary: bool = True
) -> bytes:
    """
    Generate populated report from template and data.
    
    Args:
        template_file: Template file object
        template_type: File extension
        data_df: Campaign data
        mapping_config: Field mapping configuration
        report_name: Output report name
        aggregation: Data aggregation method
        include_charts: Whether to include charts
        include_summary: Whether to include summary
    
    Returns:
        Generated report as bytes
    """
    logger.info("=" * 50)
    logger.info(f"GENERATING REPORT: {report_name}")
    logger.info(f"  Template type: {template_type}")
    logger.info(f"  Data rows: {len(data_df)}")
    logger.info(f"  Data columns: {list(data_df.columns)}")
    logger.info(f"  Mapping mode: {mapping_config.get('_mode', 'placeholder')}")
    logger.info(f"  Aggregation: {aggregation}")
    logger.info("=" * 50)
    
    if template_type in ['xlsx', 'xls']:
        return generate_excel_report(
            template_file, data_df, mapping_config, aggregation, include_charts, include_summary
        )
    elif template_type == 'csv':
        return generate_csv_report(data_df, mapping_config)
    elif template_type == 'pptx':
        return generate_pptx_report(
            template_file, data_df, mapping_config, include_charts
        )
    else:
        raise ValueError(f"Unsupported template type: {template_type}")


def generate_excel_report(
    template_file,
    data_df: pd.DataFrame,
    mapping_config: Dict,
    aggregation: str,
    include_charts: bool,
    include_summary: bool
) -> bytes:
    """
    Generate Excel report by replacing placeholders OR populating data rows.
    
    Supports two modes:
    1. Placeholder mode: Replace {{field}} placeholders with aggregated values
    2. Column mode: Populate data rows into template columns (preserving headers)
    
    Calculated Fields (auto-computed):
    - CTR: Clicks / Impressions * 100
    - CPC: Spend / Clicks
    - CPA: Spend / Conversions
    - CPM: Spend / Impressions * 1000
    - ROAS: Revenue / Spend
    - Conversion Rate: Conversions / Clicks * 100
    
    Process:
    1. Load template workbook
    2. Detect mode from mapping_config
    3. For placeholder mode: replace placeholders with aggregated data
    4. For column mode: populate data rows starting from row 2
    5. Auto-calculate derived metrics (CTR, CPC, CPA, ROAS, etc.)
    6. Preserve all formatting, formulas, and charts
    """
    import openpyxl
    from openpyxl.utils.dataframe import dataframe_to_rows
    from openpyxl.styles import Font, PatternFill, numbers
    from openpyxl.utils import get_column_letter
    
    # Define calculated fields and their formulas
    # Each entry: (field_aliases, numerator_aliases, denominator_aliases, multiplier, format)
    CALCULATED_FIELDS = {
        'ctr': {
            'aliases': ['ctr', 'click_through_rate', 'click through rate', 'clickthroughrate'],
            'numerator': ['clicks', 'click', 'total_clicks'],
            'denominator': ['impressions', 'impr', 'impression', 'total_impressions'],
            'multiplier': 100,
            'format': '0.00%',
            'formula_template': '=IF({denom}=0,0,{numer}/{denom}*100)'
        },
        'cpc': {
            'aliases': ['cpc', 'cost_per_click', 'cost per click', 'costperclick', 'avg_cpc', 'average_cpc'],
            'numerator': ['spend', 'cost', 'spend_usd', 'cost_usd', 'total_spend', 'amount_spent'],
            'denominator': ['clicks', 'click', 'total_clicks'],
            'multiplier': 1,
            'format': '$#,##0.00',
            'formula_template': '=IF({denom}=0,0,{numer}/{denom})'
        },
        'cpa': {
            'aliases': ['cpa', 'cost_per_acquisition', 'cost per acquisition', 'cost_per_conversion', 
                       'cost per conversion', 'costperacquisition', 'costperconversion', 'cpl', 'cost_per_lead'],
            'numerator': ['spend', 'cost', 'spend_usd', 'cost_usd', 'total_spend', 'amount_spent'],
            'denominator': ['conversions', 'conversion', 'conv', 'total_conversions', 'leads', 'acquisitions'],
            'multiplier': 1,
            'format': '$#,##0.00',
            'formula_template': '=IF({denom}=0,0,{numer}/{denom})'
        },
        'cpm': {
            'aliases': ['cpm', 'cost_per_mille', 'cost per mille', 'cost_per_thousand', 'costpermille'],
            'numerator': ['spend', 'cost', 'spend_usd', 'cost_usd', 'total_spend', 'amount_spent'],
            'denominator': ['impressions', 'impr', 'impression', 'total_impressions'],
            'multiplier': 1000,
            'format': '$#,##0.00',
            'formula_template': '=IF({denom}=0,0,{numer}/{denom}*1000)'
        },
        'roas': {
            'aliases': ['roas', 'return_on_ad_spend', 'return on ad spend', 'returnonadspend'],
            'numerator': ['revenue', 'revenue_usd', 'total_revenue', 'sales', 'value'],
            'denominator': ['spend', 'cost', 'spend_usd', 'cost_usd', 'total_spend', 'amount_spent'],
            'multiplier': 1,
            'format': '0.00',
            'formula_template': '=IF({denom}=0,0,{numer}/{denom})'
        },
        'conversion_rate': {
            'aliases': ['conversion_rate', 'conv_rate', 'cvr', 'cr'],
            'numerator': ['conversions', 'conversion', 'conv', 'total_conversions'],
            'denominator': ['clicks', 'click', 'total_clicks'],
            'multiplier': 100,
            'format': '0.00%',
            'formula_template': '=IF({denom}=0,0,{numer}/{denom}*100)'
        }
    }
    
    def normalize_field_name(name: str) -> str:
        """Normalize field name for matching."""
        return name.lower().replace('_', '').replace(' ', '').replace('-', '')
    
    def find_calculated_field(template_col: str) -> Optional[Dict]:
        """Check if a template column is a calculated field."""
        norm_col = normalize_field_name(template_col)
        for field_key, field_def in CALCULATED_FIELDS.items():
            for alias in field_def['aliases']:
                if normalize_field_name(alias) == norm_col:
                    return {'key': field_key, **field_def}
        return None
    
    def find_source_column(aliases: List[str], available_cols: List[str]) -> Optional[str]:
        """Find a source column matching any of the aliases."""
        for alias in aliases:
            norm_alias = normalize_field_name(alias)
            for col in available_cols:
                if normalize_field_name(col) == norm_alias:
                    return col
        # Partial match
        for alias in aliases:
            norm_alias = normalize_field_name(alias)
            for col in available_cols:
                norm_col = normalize_field_name(col)
                if norm_alias in norm_col or norm_col in norm_alias:
                    return col
        return None
    
    # Load template (preserving formatting)
    logger.debug("Loading template workbook...")
    wb = openpyxl.load_workbook(template_file)
    logger.debug(f"Template loaded. Sheets: {wb.sheetnames}")
    
    # Check if we're in column mode
    if mapping_config.get('_mode') == 'column':
        # COLUMN MODE: Populate data rows into template
        logger.info("Processing in COLUMN MODE")
        column_mappings = mapping_config.get('_column_mappings', {})
        logger.info(f"Column mappings: {len(column_mappings)} columns")
        
        for col_name, col_config in column_mappings.items():
            logger.debug(f"  Mapping: {col_name} -> {col_config.get('data_column')} (col {col_config.get('column_letter')})")
        
        if column_mappings:
            # Group mappings by sheet
            sheets_data = {}
            for template_col, config in column_mappings.items():
                sheet_name = config['sheet']
                if sheet_name not in sheets_data:
                    sheets_data[sheet_name] = {
                        'columns': [],
                        'header_row': config.get('header_row', 1)
                    }
                sheets_data[sheet_name]['columns'].append({
                    'template_col': template_col,
                    'data_col': config['data_column'],
                    'column_letter': config['column_letter']
                })
            
            logger.info(f"Processing {len(sheets_data)} sheet(s)")
            
            # Populate each sheet
            for sheet_name, sheet_config in sheets_data.items():
                logger.info(f"Processing sheet: {sheet_name}")
                if sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    col_configs = sheet_config['columns']
                    
                    # Detect header row by finding row with most filled cells
                    header_row = 1
                    max_filled = 0
                    for row_idx in range(1, min(21, sheet.max_row + 1)):
                        filled = sum(1 for cell in sheet[row_idx] if cell.value is not None)
                        if filled > max_filled:
                            max_filled = filled
                            header_row = row_idx
                    
                    data_start_row = header_row + 1
                    logger.debug(f"  Header row: {header_row}, Data starts at row: {data_start_row}")
                    
                    # Build column letter mapping for calculated fields
                    col_letter_map = {}  # data_col -> column_letter
                    for col_config in col_configs:
                        if col_config['data_col']:
                            col_letter_map[col_config['data_col'].lower()] = col_config['column_letter']
                    
                    # Also scan header row for unmapped columns (for calculated field references)
                    for cell in sheet[header_row]:
                        # Skip merged cells (use cell.column instead of cell.column_letter)
                        if not hasattr(cell, 'column') or cell.column is None:
                            continue
                        if cell.value:
                            cell_col_letter = get_column_letter(cell.column)
                            header_name = str(cell.value).lower().replace('_', '').replace(' ', '')
                            # Map common metric names to their column letters
                            for alias_group in [
                                (['impressions', 'impr'], 'impressions'),
                                (['clicks', 'click'], 'clicks'),
                                (['spend', 'cost', 'spendusd'], 'spend'),
                                (['conversions', 'conv'], 'conversions'),
                                (['revenue', 'revenueusd'], 'revenue'),
                            ]:
                                if header_name in [a.replace('_', '') for a in alias_group[0]]:
                                    col_letter_map[alias_group[1]] = cell_col_letter
                    
                    # Identify calculated field columns in template
                    calculated_cols = []
                    for cell in sheet[header_row]:
                        # Skip merged cells
                        if not hasattr(cell, 'column') or cell.column is None:
                            continue
                        if cell.value:
                            cell_col_letter = get_column_letter(cell.column)
                            calc_field = find_calculated_field(str(cell.value))
                            if calc_field:
                                # Find source columns for this calculated field
                                numer_col = find_source_column(calc_field['numerator'], list(data_df.columns))
                                denom_col = find_source_column(calc_field['denominator'], list(data_df.columns))
                                
                                calculated_cols.append({
                                    'column_letter': cell_col_letter,
                                    'field': calc_field,
                                    'numer_data_col': numer_col,
                                    'denom_data_col': denom_col,
                                    'numer_letter': col_letter_map.get(numer_col.lower() if numer_col else '', None),
                                    'denom_letter': col_letter_map.get(denom_col.lower() if denom_col else '', None),
                                })
                    
                    if calculated_cols:
                        logger.info(f"  Calculated fields detected: {len(calculated_cols)}")
                        for cc in calculated_cols:
                            logger.debug(f"    {cc['field']['key']}: {cc['numer_data_col']}/{cc['denom_data_col']} -> col {cc['column_letter']}")
                    
                    # Clear existing data rows (keep header row intact)
                    if sheet.max_row >= data_start_row:
                        rows_to_delete = sheet.max_row - data_start_row + 1
                        logger.debug(f"  Clearing {rows_to_delete} existing data rows")
                        sheet.delete_rows(data_start_row, rows_to_delete)
                    
                    # Write data rows starting after header
                    logger.info(f"  Writing {len(data_df)} data rows starting at row {data_start_row}")
                    
                    # Log the actual column configs being used
                    logger.debug(f"  Column configs for writing:")
                    for cc in col_configs:
                        logger.debug(f"    {cc['template_col']} -> {cc['data_col']} (col {cc['column_letter']})")
                    
                    # Track what's actually written
                    write_stats = {}
                    
                    for row_idx, (_, data_row) in enumerate(data_df.iterrows(), start=data_start_row):
                        # Write regular mapped columns
                        for col_config in col_configs:
                            col_letter = col_config['column_letter']
                            data_col = col_config['data_col']
                            
                            if col_letter and data_col in data_df.columns:
                                value = data_row[data_col]
                                # Handle NaN values
                                if pd.isna(value):
                                    value = ""
                                # Handle datetime
                                if hasattr(value, 'strftime'):
                                    value = value.strftime('%Y-%m-%d')
                                cell = sheet[f"{col_letter}{row_idx}"]
                                cell.value = value
                                
                                # Track first row writes for logging
                                if row_idx == data_start_row:
                                    write_stats[col_letter] = {'col': data_col, 'value': str(value)[:30]}
                            elif row_idx == data_start_row:
                                logger.warning(f"    Skipping {col_config['template_col']}: col_letter={col_letter}, data_col={data_col} in columns={data_col in data_df.columns}")
                        
                        # Write calculated fields for this row
                        for calc_col in calculated_cols:
                            col_letter = calc_col['column_letter']
                            field = calc_col['field']
                            numer_col = calc_col['numer_data_col']
                            denom_col = calc_col['denom_data_col']
                            numer_letter = calc_col['numer_letter']
                            denom_letter = calc_col['denom_letter']
                            
                            cell = sheet[f"{col_letter}{row_idx}"]
                            
                            # Option 1: Use Excel formula if we have column letters
                            if numer_letter and denom_letter:
                                formula = field['formula_template'].format(
                                    numer=f"{numer_letter}{row_idx}",
                                    denom=f"{denom_letter}{row_idx}"
                                )
                                cell.value = formula
                            # Option 2: Calculate value directly from data
                            elif numer_col and denom_col and numer_col in data_df.columns and denom_col in data_df.columns:
                                numer_val = data_row[numer_col]
                                denom_val = data_row[denom_col]
                                
                                if pd.notna(numer_val) and pd.notna(denom_val) and denom_val != 0:
                                    calc_value = (numer_val / denom_val) * field['multiplier']
                                    cell.value = round(calc_value, 4)
                                else:
                                    cell.value = 0
                            else:
                                cell.value = ""  # Cannot calculate
                    
                    # Log what was written in first row
                    if write_stats:
                        logger.debug(f"  First row writes:")
                        for col_letter, info in sorted(write_stats.items()):
                            logger.debug(f"    {col_letter}: {info['col']} = {info['value']}")
                    
                    # Auto-adjust column widths
                    for column in sheet.columns:
                        max_length = 0
                        first_cell = column[0]
                        # Skip merged cells
                        if not hasattr(first_cell, 'column') or first_cell.column is None:
                            continue
                        col_letter = get_column_letter(first_cell.column)
                        for cell in column:
                            try:
                                if cell.value and len(str(cell.value)) > max_length:
                                    max_length = len(str(cell.value))
                            except:
                                pass
                        adjusted_width = min(max_length + 2, 50)
                        sheet.column_dimensions[col_letter].width = adjusted_width
        
        # Save and return
        logger.info("Column mode processing complete. Saving workbook...")
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)
        logger.info(f"Report generated successfully ({len(output.getvalue())} bytes)")
        return output.getvalue()
    
    # PLACEHOLDER MODE: Replace placeholders with aggregated data
    logger.info("Processing in PLACEHOLDER MODE")
    for placeholder_name, config in mapping_config.items():
        if placeholder_name.startswith('_'):
            continue  # Skip internal keys
            
        data_column = config['data_column']
        location = config['location']
        
        # Parse location (e.g., "Sheet1!A1")
        if '!' in location:
            sheet_name, cell_ref = location.split('!')
            sheet = wb[sheet_name]
            
            # Calculate aggregated value based on method
            try:
                if aggregation == 'sum':
                    value = data_df[data_column].sum()
                elif aggregation == 'average':
                    value = data_df[data_column].mean()
                elif aggregation == 'latest':
                    value = data_df[data_column].iloc[-1] if len(data_df) > 0 else 0
                elif aggregation == 'all rows':
                    # For "all rows", insert comma-separated values
                    value = ', '.join(map(str, data_df[data_column].tolist()))
                else:
                    value = data_df[data_column].sum()
                
                # Format numeric values
                if isinstance(value, (int, float)):
                    value = round(value, 2)
                
                # Replace placeholder in cell
                cell = sheet[cell_ref]
                
                # If cell contains placeholder pattern, replace it
                if cell.value and isinstance(cell.value, str):
                    cell.value = cell.value.replace(placeholder_name, str(value))
                else:
                    cell.value = value
                    
            except Exception as e:
                print(f"Error processing {placeholder_name}: {e}")
                sheet[cell_ref] = f"Error: {str(e)}"
    
    # Add raw data sheet if no placeholders or if requested
    if not mapping_config or include_summary:
        # Create or get Data sheet
        if 'Campaign_Data' not in wb.sheetnames:
            data_sheet = wb.create_sheet('Campaign_Data')
        else:
            data_sheet = wb['Campaign_Data']
        
        # Clear existing data
        data_sheet.delete_rows(1, data_sheet.max_row)
        
        # Add headers with formatting
        for c_idx, col_name in enumerate(data_df.columns, 1):
            cell = data_sheet.cell(row=1, column=c_idx, value=col_name)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        
        # Add data rows
        for r_idx, row in enumerate(dataframe_to_rows(data_df, index=False, header=False), 2):
            for c_idx, value in enumerate(row, 1):
                data_sheet.cell(row=r_idx, column=c_idx, value=value)
        
        # Auto-adjust column widths
        for column in data_sheet.columns:
            max_length = 0
            first_cell = column[0]
            # Skip merged cells
            if not hasattr(first_cell, 'column') or first_cell.column is None:
                continue
            col_letter = get_column_letter(first_cell.column)
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)
            data_sheet.column_dimensions[col_letter].width = adjusted_width
    
    # Save to bytes
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    return output.getvalue()


def generate_csv_report(data_df: pd.DataFrame, mapping_config: Dict) -> bytes:
    """Generate CSV report."""
    output = io.StringIO()
    data_df.to_csv(output, index=False)
    return output.getvalue().encode('utf-8')


def generate_pptx_report(
    template_file,
    data_df: pd.DataFrame,
    mapping_config: Dict,
    include_charts: bool
) -> bytes:
    """
    Generate PowerPoint report by replacing placeholders with actual data.
    
    Process:
    1. Load template presentation
    2. Scan all slides for text placeholders
    3. Replace placeholders with aggregated data
    4. Preserve all formatting, layouts, and existing charts
    5. Optionally add data summary slide
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    # Load template (preserving all formatting and layouts)
    prs = Presentation(template_file)
    
    # Replace text placeholders in all slides
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Check text boxes and text frames
            if hasattr(shape, "text") and shape.text:
                original_text = shape.text
                updated_text = original_text
                
                # Replace all mapped placeholders
                for placeholder_name, config in mapping_config.items():
                    if placeholder_name in updated_text:
                        data_column = config['data_column']
                        
                        try:
                            # Calculate aggregated value
                            value = data_df[data_column].sum()
                            
                            # Format numeric values
                            if isinstance(value, (int, float)):
                                if value >= 1000000:
                                    formatted_value = f"${value/1000000:.2f}M"
                                elif value >= 1000:
                                    formatted_value = f"${value/1000:.1f}K"
                                else:
                                    formatted_value = f"${value:,.2f}"
                            else:
                                formatted_value = str(value)
                            
                            # Replace placeholder
                            updated_text = updated_text.replace(placeholder_name, formatted_value)
                        except Exception as e:
                            print(f"Error processing {placeholder_name} on slide {slide_idx + 1}: {e}")
                            updated_text = updated_text.replace(placeholder_name, "N/A")
                
                # Update text if changed
                if updated_text != original_text:
                    if hasattr(shape, "text_frame"):
                        shape.text_frame.text = updated_text
                    else:
                        shape.text = updated_text
            
            # Check tables for placeholders
            if hasattr(shape, "table"):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            original_text = cell.text
                            updated_text = original_text
                            
                            for placeholder_name, config in mapping_config.items():
                                if placeholder_name in updated_text:
                                    data_column = config['data_column']
                                    try:
                                        value = data_df[data_column].sum()
                                        if isinstance(value, (int, float)):
                                            value = round(value, 2)
                                        updated_text = updated_text.replace(placeholder_name, str(value))
                                    except:
                                        updated_text = updated_text.replace(placeholder_name, "N/A")
                            
                            if updated_text != original_text:
                                cell.text = updated_text
    
    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def render_settings_page():
    """Render settings page."""
    st.markdown('<div class="main-header">⚙️ Settings</div>', unsafe_allow_html=True)
    
    st.markdown("""
    ### Report Generation Settings
    
    Configure default behaviors for report generation.
    """)
    
    st.divider()
    
    st.subheader("Default Options")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.selectbox("Default Aggregation", ["Sum", "Average", "Latest", "All Rows"])
        st.checkbox("Auto-detect placeholders", value=True)
    
    with col2:
        st.selectbox("Date Format", ["YYYY-MM-DD", "DD/MM/YYYY", "MM/DD/YYYY"])
        st.checkbox("Include metadata", value=True)
    
    st.divider()
    
    st.subheader("Template Library")
    st.info("Coming soon: Save and reuse template configurations")


def main():
    """Main application entry point."""
    init_session_state()
    page = render_sidebar()
    
    if page == "upload_data_&_template":
        render_upload_page()
    elif page == "data_mapping":
        render_data_mapping_page()
    elif page == "generate_report":
        render_generate_report_page()
    elif page == "settings":
        render_settings_page()
    
    st.divider()
    st.caption("PCA Agent Reporting Module · Automated Report Generation")


if __name__ == "__main__":
    main()
