"""
Report Assembly Agent for generating PowerPoint reports.
"""
from typing import List, Dict, Any, Optional
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from loguru import logger

from ..models.campaign import (
    Campaign,
    ChannelPerformance,
    Achievement,
    ConsolidatedReport,
    ReportConfig,
    ReportTemplate
)
from ..config.settings import settings


class ReportAgent:
    """Agent for assembling and generating PowerPoint reports."""
    
    def __init__(self, output_dir: Path = None):
        """
        Initialize Report Agent.
        
        Args:
            output_dir: Directory to save generated reports
        """
        self.output_dir = output_dir or settings.report_dir
        self.output_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"Initialized ReportAgent with output_dir: {self.output_dir}")
    
    def generate_report(
        self,
        consolidated_report: ConsolidatedReport,
        config: ReportConfig = None
    ) -> str:
        """
        Generate PowerPoint report.
        
        Args:
            consolidated_report: Consolidated campaign data
            config: Report configuration
            
        Returns:
            Path to generated report file
        """
        config = config or ReportConfig()
        campaign = consolidated_report.campaign
        
        logger.info(f"Generating {config.template.value} report for campaign {campaign.campaign_id}")
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333) if config.page_size == "16:9" else Inches(10)
        prs.slide_height = Inches(7.5) if config.page_size == "16:9" else Inches(7.5)
        
        # Add slides
        if config.include_executive_summary:
            self._add_title_slide(prs, campaign, config)
            self._add_executive_summary_slide(prs, consolidated_report, config)
        
        if config.include_channel_breakdown:
            self._add_channel_overview_slide(prs, consolidated_report, config)
            for cp in consolidated_report.channel_performances:
                self._add_channel_detail_slide(prs, cp, config)
        
        if config.include_cross_channel_analysis:
            self._add_cross_channel_slide(prs, consolidated_report, config)
        
        if config.include_achievements:
            self._add_achievements_slide(prs, consolidated_report, config)
        
        if config.include_recommendations:
            self._add_recommendations_slide(prs, consolidated_report, config)
        
        if config.include_visualizations:
            self._add_visualization_slides(prs, consolidated_report, config)
        
        # Save report
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{campaign.campaign_id}_report_{timestamp}.pptx"
        filepath = self.output_dir / filename
        
        prs.save(str(filepath))
        logger.info(f"Report saved to {filepath}")
        
        return str(filepath)
    
    def _add_title_slide(
        self,
        prs: Presentation,
        campaign: Campaign,
        config: ReportConfig
    ):
        """Add title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = campaign.campaign_name
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.8), Inches(11.333), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Post Campaign Analysis Report"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Date range
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.8), Inches(11.333), Inches(0.5)
        )
        date_frame = date_box.text_frame
        date_frame.text = f"{campaign.date_range.start} to {campaign.date_range.end}"
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(18)
        date_para.alignment = PP_ALIGN.CENTER
        
        # Company name
        company_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(11.333), Inches(0.5)
        )
        company_frame = company_box.text_frame
        company_frame.text = config.company_name or settings.company_name
        company_para = company_frame.paragraphs[0]
        company_para.font.size = Pt(16)
        company_para.alignment = PP_ALIGN.CENTER
    
    def _add_executive_summary_slide(
        self,
        prs: Presentation,
        report: ConsolidatedReport,
        config: ReportConfig
    ):
        """Add executive summary slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        # Content
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(11.333), Inches(5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # Summary text
        p = tf.paragraphs[0]
        p.text = report.executive_summary
        p.font.size = Pt(16)
        p.space_after = Pt(20)
        
        # Key metrics
        p = tf.add_paragraph()
        p.text = "Key Metrics:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_before = Pt(20)
        p.space_after = Pt(10)
        
        metrics = [
            f"Total Spend: ${report.total_spend:,.2f}",
            f"Total Conversions: {report.total_conversions:,.0f}",
            f"Overall ROAS: {report.overall_roas:.2f}x" if report.overall_roas else None,
            f"Channels Analyzed: {len(report.channel_performances)}"
        ]
        
        for metric in metrics:
            if metric:
                p = tf.add_paragraph()
                p.text = f"• {metric}"
                p.font.size = Pt(14)
                p.level = 1
    
    def _add_channel_overview_slide(
        self,
        prs: Presentation,
        report: ConsolidatedReport,
        config: ReportConfig
    ):
        """Add channel overview comparison slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Channel Performance Overview"
        
        # Create table
        rows = len(report.channel_performances) + 1
        cols = 6
        
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(12.333)
        height = Inches(4.5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Header row
        headers = ["Channel", "Spend", "Conversions", "ROAS", "CPA", "Score"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        # Data rows
        for i, cp in enumerate(report.channel_performances, 1):
            table.cell(i, 0).text = cp.platform_name
            table.cell(i, 1).text = f"${cp.total_spend:,.0f}" if cp.total_spend else "N/A"
            table.cell(i, 2).text = f"{cp.total_conversions:,.0f}" if cp.total_conversions else "N/A"
            table.cell(i, 3).text = f"{cp.roas:.2f}x" if cp.roas else "N/A"
            table.cell(i, 4).text = f"${cp.cpa:.2f}" if cp.cpa else "N/A"
            table.cell(i, 5).text = f"{cp.performance_score:.1f}" if cp.performance_score else "N/A"
            
            # Set font size for data
            for j in range(cols):
                table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(11)
    
    def _add_channel_detail_slide(
        self,
        prs: Presentation,
        channel_performance: ChannelPerformance,
        config: ReportConfig
    ):
        """Add detailed slide for a single channel."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = f"{channel_performance.platform_name} Performance"
        
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(11.333), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # Performance score
        p = tf.paragraphs[0]
        p.text = f"Performance Score: {channel_performance.performance_score:.1f}/100"
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(15)
        
        # Key metrics
        p = tf.add_paragraph()
        p.text = "Key Metrics:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_after = Pt(10)
        
        metrics = [
            ("Impressions", channel_performance.total_impressions, "{:,.0f}"),
            ("Clicks", channel_performance.total_clicks, "{:,.0f}"),
            ("CTR", channel_performance.ctr, "{:.2f}%"),
            ("Conversions", channel_performance.total_conversions, "{:,.0f}"),
            ("Spend", channel_performance.total_spend, "${:,.2f}"),
            ("CPA", channel_performance.cpa, "${:.2f}"),
            ("ROAS", channel_performance.roas, "{:.2f}x"),
        ]
        
        for label, value, fmt in metrics:
            if value is not None:
                p = tf.add_paragraph()
                p.text = f"• {label}: {fmt.format(value)}"
                p.font.size = Pt(13)
                p.level = 1
        
        # Strengths
        if channel_performance.strengths:
            p = tf.add_paragraph()
            p.text = "Strengths:"
            p.font.size = Pt(16)
            p.font.bold = True
            p.space_before = Pt(15)
            p.space_after = Pt(10)
            
            for strength in channel_performance.strengths:
                p = tf.add_paragraph()
                p.text = f"✓ {strength}"
                p.font.size = Pt(12)
                p.level = 1
        
        # Opportunities
        if channel_performance.opportunities:
            p = tf.add_paragraph()
            p.text = "Opportunities:"
            p.font.size = Pt(16)
            p.font.bold = True
            p.space_before = Pt(15)
            p.space_after = Pt(10)
            
            for opp in channel_performance.opportunities:
                p = tf.add_paragraph()
                p.text = f"→ {opp}"
                p.font.size = Pt(12)
                p.level = 1
    
    def _add_cross_channel_slide(
        self,
        prs: Presentation,
        report: ConsolidatedReport,
        config: ReportConfig
    ):
        """Add cross-channel insights slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Cross-Channel Insights"
        
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(11.333), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, insight in enumerate(report.cross_channel_insights[:5], 1):
            p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
            p.text = f"{i}. {insight.title}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.space_after = Pt(8)
            
            p = tf.add_paragraph()
            p.text = insight.description
            p.font.size = Pt(13)
            p.level = 1
            p.space_after = Pt(15)
    
    def _add_achievements_slide(
        self,
        prs: Presentation,
        report: ConsolidatedReport,
        config: ReportConfig
    ):
        """Add achievements slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "🏆 Key Achievements"
        
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(11.333), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, achievement in enumerate(report.achievements[:5], 1):
            p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
            
            # Achievement icon based on impact
            icon = "🌟" if achievement.impact_level == "high" else "⭐" if achievement.impact_level == "medium" else "✨"
            
            p.text = f"{icon} {achievement.title}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.space_after = Pt(8)
            
            p = tf.add_paragraph()
            p.text = achievement.description
            p.font.size = Pt(13)
            p.level = 1
            p.space_after = Pt(15)
    
    def _add_recommendations_slide(
        self,
        prs: Presentation,
        report: ConsolidatedReport,
        config: ReportConfig
    ):
        """Add recommendations slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Strategic Recommendations"
        
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(11.333), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, recommendation in enumerate(report.recommendations, 1):
            p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
            p.text = f"{i}. {recommendation}"
            p.font.size = Pt(14)
            p.space_after = Pt(12)
    
    def _add_visualization_slides(
        self,
        prs: Presentation,
        report: ConsolidatedReport,
        config: ReportConfig
    ):
        """Add slides with visualizations."""
        for viz in report.visualizations:
            if viz.get("filepath") and Path(viz["filepath"]).exists():
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                
                title = slide.shapes.title
                title.text = viz.get("title", "Visualization")
                
                # Add image
                left = Inches(1.5)
                top = Inches(2)
                width = Inches(10)
                
                slide.shapes.add_picture(
                    viz["filepath"],
                    left, top,
                    width=width
                )
